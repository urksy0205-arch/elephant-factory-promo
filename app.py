# -*- coding: utf-8 -*-
"""
코끼리공장 다국어 홍보물 자동 생성 시스템 - PPT 편집 기능 추가
"""

import streamlit as st
from PIL import Image, ImageDraw, ImageFont
from deep_translator import GoogleTranslator
import time
from datetime import datetime
from pathlib import Path
import docx
import PyPDF2
import io
import zipfile
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ============================================
# 페이지 설정
# ============================================

st.set_page_config(
    page_title="코끼리공장 홍보물 생성기",
    page_icon="🐘",
    layout="wide"
)

# ============================================
# 설정
# ============================================

LANGUAGES = {
    'ko': '한국어 🇰🇷',
    'en': 'English 🇺🇸',
    'ja': '日本語 🇯🇵',
    'zh-CN': '中文(简体) 🇨🇳',
    'vi': 'Tiếng Việt 🇻🇳',
    'ru': 'Русский 🇷🇺',
    'uz': "O'zbek 🇺🇿",
    'si': 'සිංහල 🇱🇰'
}

BRAND_COLOR = '#2B9FD9'
BRAND_COLOR_RGB = (43, 159, 217)  # RGB 값
ACCENT_COLOR_RGB = (255, 107, 107)  # 주황색

# ============================================
# CSS 스타일
# ============================================

st.markdown("""
<style>
    .main-header {
        text-align: center;
        padding: 2rem;
        background: linear-gradient(135deg, #2B9FD9 0%, #1E88C7 100%);
        color: white;
        border-radius: 10px;
        margin-bottom: 2rem;
    }
    .stButton>button {
        width: 100%;
        background-color: #2B9FD9;
        color: white;
        font-size: 18px;
        padding: 0.5rem 1rem;
        border-radius: 5px;
        border: none;
    }
    .stButton>button:hover {
        background-color: #1E88C7;
    }
    .summary-box {
        padding: 1.5rem;
        background-color: #fff3cd;
        border-left: 4px solid #ffc107;
        border-radius: 5px;
        margin: 1rem 0;
    }
    .promo-box {
        padding: 1.5rem;
        background-color: #d4edda;
        border-left: 4px solid #28a745;
        border-radius: 5px;
        margin: 1rem 0;
    }
    .original-box {
        padding: 1.5rem;
        background-color: #f8f9fa;
        border-left: 4px solid #6c757d;
        border-radius: 5px;
        margin: 1rem 0;
    }
</style>
""", unsafe_allow_html=True)

# ============================================
# 무료 AI 요약 함수 (규칙 기반)
# ============================================

def extract_key_info(text):
    """공문에서 핵심 정보 추출"""
    info = {
        'title': '',
        'date': '',
        'time': '',
        'location': '',
        'target': '',
        'contact': '',
        'how_to_apply': '',
        'content': ''
    }
    
    lines = text.strip().split('\n')
    lines = [line.strip() for line in lines if line.strip()]
    
    # 제목 찾기
    for i, line in enumerate(lines[:5]):
        if len(line) > 5 and (
            '안내' in line or '공고' in line or '모집' in line or 
            '프로그램' in line or '교육' in line or i == 0
        ):
            info['title'] = line
            break
    
    # 날짜 찾기
    date_patterns = [
        r'(\d{4})[년.-]\s*(\d{1,2})[월.-]\s*(\d{1,2})일?',
        r'(\d{1,2})[월/]\s*(\d{1,2})일?',
        r'(\d{4})[./]\s*(\d{1,2})[./]\s*(\d{1,2})'
    ]
    
    for line in lines:
        for pattern in date_patterns:
            match = re.search(pattern, line)
            if match:
                info['date'] = match.group(0)
                break
        if info['date']:
            break
    
    # 시간 찾기
    time_patterns = [
        r'(\d{1,2}):(\d{2})',
        r'(\d{1,2})시\s*(\d{1,2})?분?'
    ]
    
    for line in lines:
        for pattern in time_patterns:
            match = re.search(pattern, line)
            if match:
                info['time'] = match.group(0)
                break
        if info['time']:
            break
    
    # 장소 찾기
    location_keywords = ['장소', '위치', '주소', '에서', '교육실', '강당']
    for line in lines:
        for keyword in location_keywords:
            if keyword in line:
                info['location'] = line
                break
        if info['location']:
            break
    
    # 대상 찾기
    target_keywords = ['대상', '참가자', '신청자', '이주민', '외국인']
    for line in lines:
        for keyword in target_keywords:
            if keyword in line:
                info['target'] = line
                break
        if info['target']:
            break
    
    # 연락처 찾기
    contact_patterns = [
        r'0\d{1,2}-\d{3,4}-\d{4}',
        r'\d{3}-\d{4}-\d{4}',
        r'010-\d{4}-\d{4}'
    ]
    
    for line in lines:
        if '연락' in line or '문의' in line or '전화' in line:
            info['contact'] = line
            for pattern in contact_patterns:
                match = re.search(pattern, line)
                if match:
                    info['contact'] = line
                    break
            break
    
    # 신청 방법 찾기
    apply_keywords = ['신청', '접수', '등록', '참여방법']
    for line in lines:
        for keyword in apply_keywords:
            if keyword in line:
                info['how_to_apply'] = line
                break
        if info['how_to_apply']:
            break
    
    # 전체 내용
    info['content'] = '\n'.join(lines)
    
    return info

def create_summary(info):
    """추출된 정보를 요약문으로 변환"""
    summary_parts = []
    
    if info['title']:
        summary_parts.append(f"📢 {info['title']}")
    
    if info['date']:
        summary_parts.append(f"📅 일시: {info['date']}")
    
    if info['time']:
        if not info['date']:
            summary_parts.append(f"🕐 시간: {info['time']}")
        else:
            summary_parts[-1] += f" {info['time']}"
    
    if info['location']:
        summary_parts.append(f"📍 {info['location']}")
    
    if info['target']:
        summary_parts.append(f"👥 {info['target']}")
    
    if info['how_to_apply']:
        summary_parts.append(f"✍️ {info['how_to_apply']}")
    
    if info['contact']:
        summary_parts.append(f"📞 {info['contact']}")
    
    return '\n'.join(summary_parts)

def create_promo_text(info):
    """홍보문 스타일로 변환"""
    promo_parts = []
    
    # 제목
    if info['title']:
        title = info['title'].replace('안내', '').replace('공고', '').strip()
        promo_parts.append(f"🎉 {title} 🎉")
    else:
        promo_parts.append("🎉 코끼리공장에서 알려드립니다! 🎉")
    
    promo_parts.append("")
    
    # 핵심 내용
    content_line = "코끼리공장에서 이주민 여러분을 위한 프로그램을 준비했습니다! 💙"
    
    if '교육' in info['content']:
        content_line = "이주민을 위한 무료 교육 프로그램에 참여하세요! 📚"
    elif '모집' in info['content']:
        content_line = "여러분의 참여를 기다립니다! 함께해요! 🙌"
    elif '행사' in info['content']:
        content_line = "즐거운 행사에 여러분을 초대합니다! 🎊"
    
    promo_parts.append(content_line)
    promo_parts.append("")
    
    # 핵심 정보
    if info['date'] or info['time']:
        date_str = info['date'] if info['date'] else ''
        time_str = info['time'] if info['time'] else ''
        promo_parts.append(f"📅 {date_str} {time_str}".strip())
    
    if info['location']:
        location = info['location'].replace('장소:', '').replace('장소', '').strip()
        promo_parts.append(f"📍 {location}")
    
    promo_parts.append("")
    
    # 참여 유도
    if info['how_to_apply']:
        apply = info['how_to_apply'].replace('신청:', '').replace('신청', '').strip()
        promo_parts.append(f"✅ {apply}")
    else:
        promo_parts.append("✅ 지금 바로 신청하세요!")
    
    if info['contact']:
        promo_parts.append(f"📞 {info['contact']}")
    
    promo_parts.append("")
    promo_parts.append("💙 많은 참여 바랍니다! 💙")
    
    return '\n'.join(promo_parts)

# ============================================
# 파일 읽기 함수
# ============================================

def read_docx(file):
    """워드 파일 읽기"""
    doc = docx.Document(file)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return '\n'.join(text)

def read_pdf(file):
    """PDF 파일 읽기"""
    text = []
    pdf = PyPDF2.PdfReader(file)
    for page in pdf.pages:
        text.append(page.extract_text())
    return '\n'.join(text)

def read_txt(file):
    """텍스트 파일 읽기"""
    return file.read().decode('utf-8')

def translate_text(text, target_lang):
    """텍스트 번역"""
    try:
        if target_lang == 'ko':
            return text
        
        translator = GoogleTranslator(source='ko', target=target_lang)
        return translator.translate(text)
    except Exception as e:
        st.warning(f"번역 실패 ({target_lang}): {str(e)}")
        return text

# ============================================
# PPT 생성 함수 (NEW!)
# ============================================

def create_ppt_slide(title, content, lang_code, size_type='social'):
    """편집 가능한 PPT 슬라이드 생성"""
    
    # 프레젠테이션 생성
    prs = Presentation()
    
    # 슬라이드 크기 설정
    if size_type == 'social':
        # 소셜미디어용 (정사각형)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(10)
    else:  # a4
        # A4 세로
        prs.slide_width = Inches(8.27)
        prs.slide_height = Inches(11.69)
    
    # 빈 슬라이드 레이아웃 사용
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 배경 - 흰색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 상단 파란색 바
    header_height = prs.slide_height * 0.15
    header = slide.shapes.add_shape(
        1,  # 직사각형
        Inches(0), Inches(0),
        prs.slide_width, header_height
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*BRAND_COLOR_RGB)
    header.line.fill.background()
    
    # 하단 주황색 바
    footer_height = prs.slide_height * 0.05
    footer = slide.shapes.add_shape(
        1,  # 직사각형
        Inches(0), prs.slide_height - footer_height,
        prs.slide_width, footer_height
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(*ACCENT_COLOR_RGB)
    footer.line.fill.background()
    
    # 로고 추가 (있는 경우)
    try:
        if Path('logos/logo.png').exists():
            logo_left = Inches(0.3)
            logo_top = Inches(0.3)
            logo_width = prs.slide_width * 0.25
            
            slide.shapes.add_picture(
                'logos/logo.png',
                logo_left, logo_top,
                width=logo_width
            )
    except Exception as e:
        pass
    
    # 제목 텍스트 박스
    title_left = Inches(0.5)
    title_top = prs.slide_height * 0.25
    title_width = prs.slide_width - Inches(1)
    title_height = Inches(1)
    
    title_box = slide.shapes.add_textbox(
        title_left, title_top, title_width, title_height
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    # 제목 텍스트 설정
    title_clean = re.sub(r'[🎉🎊💙❤️✨🌟⭐]', '', title).strip()
    p = title_frame.paragraphs[0]
    p.text = title_clean
    p.font.size = Pt(44 if size_type == 'social' else 54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.alignment = PP_ALIGN.CENTER
    
    # 내용 텍스트 박스
    content_left = Inches(0.5)
    content_top = prs.slide_height * 0.4
    content_width = prs.slide_width - Inches(1)
    content_height = prs.slide_height * 0.5
    
    content_box = slide.shapes.add_textbox(
        content_left, content_top, content_width, content_height
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # 내용 텍스트 설정
    lines = content.split('\n')
    for i, line in enumerate(lines[:10]):  # 최대 10줄
        if i > 0:
            content_frame.add_paragraph()
        
        p = content_frame.paragraphs[i]
        
        # 이모지 제거하고 텍스트만
        line_clean = line.strip()
        
        # 특정 이모지는 유지하고 싶다면:
        # line_clean = line.strip()
        
        p.text = line_clean
        p.font.size = Pt(20 if size_type == 'social' else 24)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(10)
        
        # 강조 표시 (📅, 📍, 📞 등이 있는 줄)
        if any(emoji in line for emoji in ['📅', '📍', '📞', '✅']):
            p.font.bold = True
    
    # PPT를 바이트로 변환
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    
    return ppt_bytes.getvalue()

# ============================================
# 이미지 생성 함수 (수정됨)
# ============================================

def create_promo_image(title, content, lang_code, size_type='social'):
    """홍보 이미지 생성"""
    
    # 크기 설정
    if size_type == 'social':
        width, height = 1080, 1080
    else:  # a4
        width, height = 2480, 3508
    
    # 배경 생성
    img = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(img)
    
    # 상단 파란색 바
    header_height = int(height * 0.15)
    draw.rectangle([(0, 0), (width, header_height)], fill=BRAND_COLOR)
    
    # 하단 주황색 바
    footer_height = int(height * 0.05)
    draw.rectangle(
        [(0, height - footer_height), (width, height)], 
        fill='#FF6B6B'
    )
    
    # 로고 추가 (있는 경우)
    try:
        if Path('logos/logo.png').exists():
            logo = Image.open('logos/logo.png')
            logo_width = int(width * 0.3)
            logo_height = int(logo_width * logo.size[1] / logo.size[0])
            logo = logo.resize((logo_width, logo_height), Image.Resampling.LANCZOS)
            
            if logo.mode != 'RGBA':
                logo = logo.convert('RGBA')
            
            img.paste(logo, (30, 30), logo)
    except:
        pass
    
    # 폰트 설정 (한글 지원 폰트 사용)
    try:
        # Windows
        title_font = ImageFont.truetype("malgun.ttf", int(height * 0.05))
        content_font = ImageFont.truetype("malgun.ttf", int(height * 0.03))
        emoji_font = ImageFont.truetype("seguiemj.ttf", int(height * 0.03))
    except:
        try:
            # Mac
            title_font = ImageFont.truetype("/System/Library/Fonts/AppleSDGothicNeo.ttc", int(height * 0.05))
            content_font = ImageFont.truetype("/System/Library/Fonts/AppleSDGothicNeo.ttc", int(height * 0.03))
            emoji_font = content_font
        except:
            try:
                # Linux
                title_font = ImageFont.truetype("/usr/share/fonts/truetype/nanum/NanumGothic.ttf", int(height * 0.05))
                content_font = ImageFont.truetype("/usr/share/fonts/truetype/nanum/NanumGothic.ttf", int(height * 0.03))
                emoji_font = content_font
            except:
                # 기본 폰트 (최후의 수단)
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
                emoji_font = content_font
    
    # 제목 그리기 (이모지 제거)
    title_y = int(height * 0.2)
    title_clean = re.sub(r'[^\w\s가-힣]', '', title).strip()
    
    # 제목을 중앙 정렬로 그리기
    title_bbox = draw.textbbox((0, 0), title_clean[:50], font=title_font)
    title_width = title_bbox[2] - title_bbox[0]
    title_x = (width - title_width) // 2
    
    draw.text((title_x, title_y), title_clean[:50], fill='#333333', font=title_font)
    
    # 내용 그리기 (카드뉴스 스타일)
    content_y = int(height * 0.35)
    line_height = int(height * 0.06)
    
    lines = content.split('\n')
    
    # 박스 스타일로 각 줄 그리기
    y_position = content_y
    
    for i, line in enumerate(lines[:10]):  # 최대 10줄
        line = line.strip()
        if not line:
            continue
        
        # 이모지와 텍스트 분리
        emoji_match = re.match(r'^([📅📍📞✅💙🎉🎊📚🙌✨]+)\s*(.+)$', line)
        
        if emoji_match:
            emoji = emoji_match.group(1)
            text = emoji_match.group(2)
            
            # 배경 박스 그리기 (연한 회색)
            box_padding = 20
            text_bbox = draw.textbbox((0, 0), text, font=content_font)
            text_width = text_bbox[2] - text_bbox[0]
            
            box_x1 = 50
            box_y1 = y_position - 10
            box_x2 = width - 50
            box_y2 = y_position + line_height - 10
            
            # 중요 정보는 노란색 박스
            if any(e in emoji for e in ['📅', '📍', '📞']):
                box_color = '#FFF9E6'
                border_color = '#FFD700'
            else:
                box_color = '#F5F5F5'
                border_color = '#DDDDDD'
            
            # 박스 그리기
            draw.rectangle([box_x1, box_y1, box_x2, box_y2], fill=box_color, outline=border_color, width=2)
            
            # 이모지 그리기
            try:
                draw.text((box_x1 + 15, y_position), emoji, fill='#333333', font=emoji_font, embedded_color=True)
            except:
                pass
            
            # 텍스트 그리기
            draw.text((box_x1 + 60, y_position), text[:50], fill='#333333', font=content_font)
            
        else:
            # 일반 텍스트 (이모지 없음)
            text_clean = re.sub(r'[^\w\s가-힣:/-]', '', line)
            draw.text((70, y_position), text_clean[:60], fill='#333333', font=content_font)
        
        y_position += line_height
    
    return img


# ============================================
# 메인 UI
# ============================================

# 헤더
st.markdown("""
<div class="main-header">
    <h1>🐘 코끼리공장 다국어 홍보물 자동 생성기</h1>
    <p>공문을 자동으로 요약하고 홍보문으로 변환한 후 8개 언어로 번역합니다</p>
    <p style="font-size: 14px; margin-top: 10px;">✨ 완전 무료 | AI 자동 요약 | 다국어 번역 | 이미지 & PPT 생성</p>
</div>
""", unsafe_allow_html=True)

# 사이드바
with st.sidebar:
    st.header("⚙️ 설정")
    
    logo_file = st.file_uploader(
        "로고 업로드 (선택사항)",
        type=['png', 'jpg', 'jpeg'],
        help="홍보물에 들어갈 로고를 업로드하세요"
    )
    
    if logo_file:
        Path('logos').mkdir(exist_ok=True)
        with open('logos/logo.png', 'wb') as f:
            f.write(logo_file.read())
        st.success("✅ 로고 업로드 완료!")
    
    st.markdown("---")
    
    st.markdown("""
    ### ✨ 새로운 기능!
    - 🎨 **PPT 편집 기능** (NEW!)
    - 🤖 AI 자동 요약
    - 📝 홍보문 자동 생성
    - 🌏 8개 언어 번역
    - 🖼️ 이미지 자동 생성
    
    ### 📋 지원 파일
    - Word (.docx)
    - PDF (.pdf)
    - Text (.txt)
    
    ### 💡 작동 방식
    1. 공문 업로드
    2. AI가 핵심 정보 추출
    3. 홍보문 스타일로 변환
    4. 다국어 번역
    5. 이미지 & PPT 생성
    """)

# 메인 영역
tab1, tab2, tab3 = st.tabs(["📝 공문 입력 & 생성", "💡 예시 보기", "ℹ️ 사용 방법"])

with tab1:
    st.header("1️⃣ 공문 입력")
    
    # 입력 방식 선택
    input_method = st.radio(
        "입력 방식을 선택하세요:",
        ["📁 파일 업로드", "✏️ 직접 입력"],
        horizontal=True
    )
    
    text_content = None
    
    if input_method == "📁 파일 업로드":
        uploaded_file = st.file_uploader(
            "공문 파일을 선택하세요",
            type=['docx', 'pdf', 'txt'],
            help="워드, PDF, 텍스트 파일을 지원합니다"
        )
        
        if uploaded_file:
            with st.spinner("파일을 읽는 중..."):
                try:
                    if uploaded_file.name.endswith('.docx'):
                        text_content = read_docx(uploaded_file)
                    elif uploaded_file.name.endswith('.pdf'):
                        text_content = read_pdf(uploaded_file)
                    elif uploaded_file.name.endswith('.txt'):
                        text_content = read_txt(uploaded_file)
                    
                    st.success(f"✅ 파일 읽기 완료! ({len(text_content)}자)")
                
                except Exception as e:
                    st.error(f"❌ 파일 읽기 실패: {str(e)}")
    
    else:  # 직접 입력
        text_content = st.text_area(
            "공문 내용을 입력하세요",
            height=300,
            placeholder="""예시:

이주민 한국어 교육 프로그램 안내

일시: 2025년 1월 15일 오후 2시
장소: 코끼리공장 교육실
대상: 이주민 누구나
신청: 전화 또는 방문 접수

코끼리공장에서 이주민을 위한 무료 한국어 교육을 진행합니다.
기초부터 차근차근 배울 수 있습니다.

문의: 052-123-4567
""",
            help="Ctrl+V로 붙여넣기 가능합니다"
        )
    
    # 원문 표시
    if text_content and len(text_content) > 10:
        with st.expander("📄 원문 보기"):
            st.markdown(f'<div class="original-box">{text_content}</div>', unsafe_allow_html=True)
        
        # AI 요약 버튼
        st.markdown("---")
        st.header("2️⃣ AI 자동 요약 및 홍보문 생성")
        
        col1, col2 = st.columns([3, 1])
        
        with col1:
            st.info("💡 공문을 분석하여 핵심 정보를 추출하고 홍보문으로 변환합니다")
        
        with col2:
            analyze_button = st.button("🤖 분석 시작", type="primary", use_container_width=True)
        
        if analyze_button:
            with st.spinner("🤖 AI가 공문을 분석하고 있습니다..."):
                # 정보 추출
                info = extract_key_info(text_content)
                
                # 요약 생성
                summary = create_summary(info)
                
                # 홍보문 생성
                promo = create_promo_text(info)
                
                # 세션에 저장
                st.session_state['original'] = text_content
                st.session_state['summary'] = summary
                st.session_state['promo'] = promo
                st.session_state['info'] = info
            
            st.success("✅ 분석 완료!")
    
    # 분석 결과 표시
    if 'promo' in st.session_state:
        st.markdown("---")
        st.header("📊 분석 결과")
        
        # 요약
        st.subheader("📌 핵심 요약")
        st.markdown(f'<div class="summary-box">{st.session_state["summary"]}</div>', unsafe_allow_html=True)
        
        # 홍보문
        st.subheader("✨ 생성된 홍보문")
        
        # 편집 가능하게
        edited_promo = st.text_area(
            "홍보문 (수정 가능)",
            value=st.session_state['promo'],
            height=300,
            help="생성된 홍보문을 수정할 수 있습니다"
        )
        
        st.session_state['promo'] = edited_promo
        
        st.markdown(f'<div class="promo-box">{edited_promo}</div>', unsafe_allow_html=True)
        
        # 언어 선택
        st.markdown("---")
        st.header("3️⃣ 번역 언어 선택")
        
        col1, col2 = st.columns(2)
        
        selected_langs = []
        lang_list = list(LANGUAGES.items())
        
        with col1:
            for i in range(0, len(lang_list), 2):
                lang_code, lang_name = lang_list[i]
                if st.checkbox(lang_name, value=True, key=f"lang_{lang_code}"):
                    selected_langs.append(lang_code)
        
        with col2:
            for i in range(1, len(lang_list), 2):
                if i < len(lang_list):
                    lang_code, lang_name = lang_list[i]
                    if st.checkbox(lang_name, value=True, key=f"lang_{lang_code}"):
                        selected_langs.append(lang_code)
        
        # 출력 형식 선택 (NEW!)
        st.markdown("---")
        st.header("4️⃣ 출력 형식 선택")
        
        output_formats = st.multiselect(
            "생성할 형식을 선택하세요",
            ["🖼️ 이미지 (PNG) - 완성본", "📊 PowerPoint (PPTX) - 편집 가능"],
            default=["🖼️ 이미지 (PNG) - 완성본", "📊 PowerPoint (PPTX) - 편집 가능"],
            help="이미지는 바로 사용 가능하고, PPT는 PowerPoint나 Google Slides에서 자유롭게 편집 가능합니다"
        )
        
        # 이미지 크기 선택 (이미지 형식 선택 시만)
        size_options = []
        if "🖼️ 이미지 (PNG) - 완성본" in output_formats:
            st.subheader("이미지 크기")
            size_options = st.multiselect(
                "생성할 이미지 크기를 선택하세요",
                ["소셜미디어용 (1080x1080)", "A4 인쇄용 (2480x3508)"],
                default=["소셜미디어용 (1080x1080)", "A4 인쇄용 (2480x3508)"]
            )
        
        # PPT 크기 선택 (PPT 형식 선택 시만)
        ppt_size_options = []
        if "📊 PowerPoint (PPTX) - 편집 가능" in output_formats:
            st.subheader("PPT 크기")
            ppt_size_options = st.multiselect(
                "생성할 PPT 크기를 선택하세요",
                ["소셜미디어용 (정사각형)", "A4 인쇄용 (세로)"],
                default=["소셜미디어용 (정사각형)", "A4 인쇄용 (세로)"]
            )
        
        # 생성 버튼
        st.markdown("---")
        st.header("5️⃣ 최종 생성")
        
        if st.button("🚀 번역 및 홍보물 생성 시작!", type="primary", use_container_width=True):
            
            if not selected_langs:
                st.error("❌ 번역할 언어를 최소 1개 이상 선택해주세요")
            elif not output_formats:
                st.error("❌ 출력 형식을 최소 1개 이상 선택해주세요")
            elif "🖼️ 이미지 (PNG) - 완성본" in output_formats and not size_options:
                st.error("❌ 이미지 크기를 최소 1개 이상 선택해주세요")
            elif "📊 PowerPoint (PPTX) - 편집 가능" in output_formats and not ppt_size_options:
                st.error("❌ PPT 크기를 최소 1개 이상 선택해주세요")
            else:
                # 진행 상황 표시
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                # 결과 저장용
                translations = {}
                images = {}
                ppts = {}
                
                # 총 단계 계산
                total_steps = len(selected_langs) * (
                    1 + 
                    (len(size_options) if "🖼️ 이미지 (PNG) - 완성본" in output_formats else 0) +
                    (len(ppt_size_options) if "📊 PowerPoint (PPTX) - 편집 가능" in output_formats else 0)
                )
                current_step = 0
                
                # 번역
                status_text.text("🌏 번역 중...")
                
                for lang_code in selected_langs:
                    lang_name = LANGUAGES[lang_code]
                    status_text.text(f"🌏 번역 중... {lang_name}")
                    
                    translated = translate_text(edited_promo, lang_code)
                    translations[lang_code] = translated
                    
                    current_step += 1
                    progress_bar.progress(current_step / total_steps)
                    
                    time.sleep(0.5)
                
                # 이미지 생성
                if "🖼️ 이미지 (PNG) - 완성본" in output_formats:
                    status_text.text("🎨 이미지 생성 중...")
                    
                    for lang_code, translated_text in translations.items():
                        lang_name = LANGUAGES[lang_code]
                        
                        # 제목과 내용 분리
                        lines = translated_text.split('\n')
                        title = lines[0][:100] if lines else "공지사항"
                        content = '\n'.join(lines[1:]) if len(lines) > 1 else translated_text
                        
                        images[lang_code] = {}
                        
                        for size_option in size_options:
                            if "소셜" in size_option:
                                size_type = 'social'
                                size_name = '소셜미디어'
                            else:
                                size_type = 'a4'
                                size_name = 'A4'
                            
                            status_text.text(f"🎨 이미지 생성 중... {lang_name} ({size_name})")
                            
                            try:
                                img = create_promo_image(title, content, lang_code, size_type)
                                
                                # 이미지를 바이트로 변환
                                img_byte_arr = io.BytesIO()
                                img.save(img_byte_arr, format='PNG')
                                img_byte_arr.seek(0)
                                
                                images[lang_code][size_type] = img_byte_arr.getvalue()
                                
                            except Exception as e:
                                st.warning(f"⚠️ {lang_name} {size_name} 이미지 생성 실패: {str(e)}")
                            
                            current_step += 1
                            progress_bar.progress(current_step / total_steps)
                
                # PPT 생성 (NEW!)
                if "📊 PowerPoint (PPTX) - 편집 가능" in output_formats:
                    status_text.text("📊 PPT 생성 중...")
                    
                    for lang_code, translated_text in translations.items():
                        lang_name = LANGUAGES[lang_code]
                        
                        # 제목과 내용 분리
                        lines = translated_text.split('\n')
                        title = lines[0][:100] if lines else "공지사항"
                        content = '\n'.join(lines[1:]) if len(lines) > 1 else translated_text
                        
                        ppts[lang_code] = {}
                        
                        for ppt_size_option in ppt_size_options:
                            if "소셜" in ppt_size_option:
                                size_type = 'social'
                                size_name = '소셜미디어'
                            else:
                                size_type = 'a4'
                                size_name = 'A4'
                            
                            status_text.text(f"📊 PPT 생성 중... {lang_name} ({size_name})")
                            
                            try:
                                ppt_bytes = create_ppt_slide(title, content, lang_code, size_type)
                                ppts[lang_code][size_type] = ppt_bytes
                                
                            except Exception as e:
                                st.warning(f"⚠️ {lang_name} {size_name} PPT 생성 실패: {str(e)}")
                            
                            current_step += 1
                            progress_bar.progress(current_step / total_steps)
                
                progress_bar.progress(1.0)
                status_text.text("✅ 완료!")
                
                # 결과 표시
                st.success("🎉 홍보물 생성 완료!")
                
                st.markdown("---")
                st.header("📥 결과물 다운로드")
                
                # 탭으로 언어별 표시
                lang_tabs = st.tabs([LANGUAGES[lang] for lang in selected_langs])
                
                for idx, lang_code in enumerate(selected_langs):
                    with lang_tabs[idx]:
                        st.subheader(f"📝 번역문")
                        st.text_area(
                            f"{LANGUAGES[lang_code]} 번역 결과",
                            translations[lang_code],
                            height=200,
                            key=f"trans_{lang_code}"
                        )
                        
                        # 이미지 표시
                        if "🖼️ 이미지 (PNG) - 완성본" in output_formats and images.get(lang_code):
                            st.subheader("🖼️ 이미지 (완성본)")
                            
                            cols = st.columns(len(size_options))
                            
                            for col_idx, size_option in enumerate(size_options):
                                size_type = 'social' if "소셜" in size_option else 'a4'
                                size_name = '소셜미디어' if size_type == 'social' else 'A4'
                                
                                with cols[col_idx]:
                                    if size_type in images.get(lang_code, {}):
                                        img_bytes = images[lang_code][size_type]
                                        st.image(img_bytes, caption=f"{size_name}용", use_container_width=True)
                                        
                                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                                        filename = f"홍보물_이미지_{lang_code}_{size_type}_{timestamp}.png"
                                        
                                        st.download_button(
                                            label=f"💾 {size_name}용 다운로드",
                                            data=img_bytes,
                                            file_name=filename,
                                            mime="image/png",
                                            key=f"dl_img_{lang_code}_{size_type}"
                                        )
                        
                        # PPT 다운로드
                        if "📊 PowerPoint (PPTX) - 편집 가능" in output_formats and ppts.get(lang_code):
                            st.subheader("📊 PowerPoint (편집 가능)")
                            
                            cols = st.columns(len(ppt_size_options))
                            
                            for col_idx, ppt_size_option in enumerate(ppt_size_options):
                                size_type = 'social' if "소셜" in ppt_size_option else 'a4'
                                size_name = '소셜미디어' if size_type == 'social' else 'A4'
                                
                                with cols[col_idx]:
                                    if size_type in ppts.get(lang_code, {}):
                                        ppt_bytes = ppts[lang_code][size_type]
                                        
                                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                                        filename = f"홍보물_PPT_{lang_code}_{size_type}_{timestamp}.pptx"
                                        
                                        st.download_button(
                                            label=f"📊 {size_name}용 PPT 다운로드",
                                            data=ppt_bytes,
                                            file_name=filename,
                                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                            key=f"dl_ppt_{lang_code}_{size_type}",
                                            help="PowerPoint나 Google Slides에서 열어 자유롭게 편집하세요!"
                                        )
                                        
                                        st.info(f"✏️ {size_name}용 PPT를 다운로드하여 PowerPoint나 Google Slides에서 편집하세요!")
                
                # 일괄 다운로드
                st.markdown("---")
                st.subheader("📦 전체 다운로드")
                
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                    
                    # 원문 저장
                    zip_file.writestr("원문.txt", st.session_state['original'].encode('utf-8'))
                    
                    # 요약 저장
                    zip_file.writestr("요약.txt", st.session_state['summary'].encode('utf-8'))
                    
                    # 홍보문 저장
                    zip_file.writestr("홍보문_한국어.txt", edited_promo.encode('utf-8'))
                    
                    # 번역문 저장
                    for lang_code, text in translations.items():
                        if lang_code != 'ko':
                            filename = f"번역문/홍보문_{lang_code}.txt"
                            zip_file.writestr(filename, text.encode('utf-8'))
                    
                    # 이미지 저장
                    if images:
                        for lang_code, size_dict in images.items():
                            for size_type, img_bytes in size_dict.items():
                                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                                filename = f"이미지/홍보물_{lang_code}_{size_type}_{timestamp}.png"
                                zip_file.writestr(filename, img_bytes)
                    
                    # PPT 저장
                    if ppts:
                        for lang_code, size_dict in ppts.items():
                            for size_type, ppt_bytes in size_dict.items():
                                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                                filename = f"PPT/홍보물_{lang_code}_{size_type}_{timestamp}.pptx"
                                zip_file.writestr(filename, ppt_bytes)
                
                zip_buffer.seek(0)
                
                st.download_button(
                    label="📦 전체 파일 다운로드 (ZIP)",
                    data=zip_buffer,
                    file_name=f"코끼리공장_홍보물_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
                    mime="application/zip"
                )

with tab2:
    st.header("💡 변환 예시")
    
    st.markdown("""
    ### 공문 → 홍보문 변환 예시
    
    AI가 어떻게 변환하는지 예시를 보여드립니다.
    """)
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("📄 원본 공문")
        st.markdown("""
        ```
        이주민 한국어 교육 프로그램 운영 안내
        
        1. 목적: 이주민의 한국어 능력 향상
        2. 일시: 2025년 1월 15일(수) 14:00
        3. 장소: 코끼리공장 2층 교육실
        4. 대상: 울산 거주 이주민
        5. 내용: 기초 한국어 회화 교육
        6. 신청: 방문 또는 전화 접수
        7. 문의: 052-123-4567
        ```
        """)
    
    with col2:
        st.subheader("✨ 생성된 홍보문")
        st.markdown("""
        ```
        🎉 이주민 한국어 교육 프로그램 🎉
        
        이주민을 위한 무료 교육 프로그램에 
        참여하세요! 📚
        
        📅 2025년 1월 15일(수) 14:00
        📍 코끼리공장 2층 교육실
        
        ✅ 방문 또는 전화로 신청하세요!
        📞 문의: 052-123-4567
        
        💙 많은 참여 바랍니다! 💙
        ```
        """)
    
    st.markdown("---")
    
    st.info("""
    💡 **변환 특징**
    - 복잡한 공문 형식 → 간결하고 친근한 홍보문
    - 핵심 정보만 추출 (날짜, 장소, 신청 방법)
    - 이모지 추가로 시각적 효과
    - 참여를 유도하는 문구 포함
    """)
    
    st.markdown("---")
    st.subheader("🎨 출력 형식 비교")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("""
        ### 🖼️ 이미지 (PNG)
        
        **장점:**
        - 바로 사용 가능
        - 디자인 완성본
        - SNS 업로드 즉시 가능
        
        **단점:**
        - 수정 불가능
        - 텍스트 변경 어려움
        
        **추천 용도:**
        - 인스타그램, 페이스북 게시
        - 빠른 공유가 필요할 때
        """)
    
    with col2:
        st.markdown("""
        ### 📊 PowerPoint (PPTX)
        
        **장점:**
        - 자유롭게 편집 가능
        - 텍스트, 이미지 수정 가능
        - 색상, 위치 조정 가능
        
        **단점:**
        - 편집 프로그램 필요
        - 약간의 편집 시간 필요
        
        **추천 용도:**
        - 내용 수정이 필요할 때
        - 브랜드 컬러 변경 시
        - 여러 버전 제작 시
        """)

with tab3:
    st.header("📖 사용 방법")
    
    st.markdown("""
    ### 🚀 전체 프로세스
    
    #### 1️⃣ 공문 입력
    - 파일 업로드 (워드, PDF, 텍스트)
    - 또는 직접 복사 & 붙여넣기
    
    #### 2️⃣ AI 분석
    - "분석 시작" 버튼 클릭
    - AI가 자동으로:
      - 제목, 날짜, 장소, 연락처 등 추출
      - 핵심 내용 요약
      - 홍보문 스타일로 변환
    
    #### 3️⃣ 홍보문 수정 (선택)
    - 생성된 홍보문을 확인
    - 필요시 직접 수정 가능
    
    #### 4️⃣ 언어 선택
    - 번역할 언어 체크
    - 여러 개 동시 선택 가능
    
    #### 5️⃣ 출력 형식 선택 ⭐ NEW!
    - **이미지 (PNG)**: 완성본, 바로 사용
    - **PowerPoint (PPTX)**: 편집 가능, 자유로운 수정
    - 둘 다 선택 가능!
    
    #### 6️⃣ 크기 선택
    - 소셜미디어용 (정사각형)
    - A4 인쇄용 (세로)
    
    #### 7️⃣ 생성 & 다운로드
    - "생성 시작" 버튼 클릭
    - 자동으로 번역 및 파일 생성
    - 개별 또는 일괄 다운로드
    
    ---
    
    ### 📊 PPT 편집 방법
    
    #### PowerPoint에서 편집하기:
    1. 다운로드한 PPTX 파일 열기
    2. 텍스트 클릭하여 직접 수정
    3. 색상, 폰트, 크기 변경 가능
    4. 이미지 추가/삭제 가능
    5. PNG로 내보내기 (파일 → 다른 이름으로 저장 → PNG)
    
    #### Google Slides에서 편집하기:
    1. Google Drive에 PPTX 파일 업로드
    2. 파일 우클릭 → Google Slides로 열기
    3. 온라인에서 바로 편집
    4. 파일 → 다운로드 → PNG 이미지
    
    ---
    
    ### 🌏 지원 언어
    
    - 🇰🇷 한국어
    - 🇺🇸 영어
    - 🇯🇵 일본어
    - 🇨🇳 중국어(간체)
    - 🇻🇳 베트남어
    - 🇷🇺 러시아어
    - 🇺🇿 우즈베키스탄어
    - 🇱🇰 스리랑카어
    
    ---
    
    ### 💡 팁
    
    1. **공문 작성 팁**
       - 날짜, 시간, 장소를 명확히 표기
       - 연락처 포함
       - 신청 방법 명시
    
    2. **더 좋은 결과를 위해**
       - 공문이 너무 길면 핵심만 입력
       - 중요한 정보는 앞부분에 배치
       - 생성 후 홍보문을 검토하고 수정
    
    3. **형식 선택 가이드**
       - 빠른 공유 필요 → 이미지 (PNG)
       - 내용 수정 필요 → PowerPoint (PPTX)
       - 확실하지 않으면 → 둘 다 생성!
    
    4. **이미지 활용**
       - 소셜미디어: 인스타그램, 페이스북
       - A4: 포스터, 전단지 인쇄
    
    ---
    
    ### ⚠️ 주의사항
    
    - ✅ 완전 무료로 사용 가능
    - ✅ 인터넷 연결 필요 (번역 기능)
    - ✅ 한글 파일(.hwp)은 미지원
    - ✅ 생성된 홍보문은 반드시 검토 후 사용
    - ✅ PPT 파일은 PowerPoint 2007 이상 또는 Google Slides에서 열 수 있습니다
    
    ---
    
    ### 📞 문의
    
    울산 코끼리공장  
    [연락처 입력]
    """)

# 푸터
st.markdown("---")
st.markdown("""
<div style="text-align: center; color: #666; padding: 1rem;">
    🐘 코끼리공장 다국어 홍보물 자동 생성기 v2.1<br>
    ✨ PPT 편집 기능 추가 | AI 자동 요약 | 완전 무료<br>
    Made with ❤️ for Elephant Factory
</div>
""", unsafe_allow_html=True)

